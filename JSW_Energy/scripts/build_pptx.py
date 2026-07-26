"""Companion pitch deck for the JSW Energy note. 18 slides, one idea per slide,
max 6 table columns per slide (oldest years dropped first)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import jswdata as d

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
INK = RGBColor(0x0B, 0x0B, 0x0B)
GREY = RGBColor(0x52, 0x51, 0x4E)
LGREY = RGBColor(0x8A, 0x88, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xEE, 0xF3, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5
M = 0.62


def slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(M), Inches(0.34), Inches(W - 2 * M), Inches(0.62))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    r.font.name = "Calibri"
    if kicker:
        tb2 = s.shapes.add_textbox(Inches(M), Inches(0.94), Inches(W - 2 * M),
                                   Inches(0.44))
        p = tb2.text_frame.paragraphs[0]
        p.text_frame = None
        r = p.add_run(); r.text = kicker
        r.font.size = Pt(14); r.font.color.rgb = ORANGE; r.font.name = "Calibri"
        r.font.bold = True
    ln = s.shapes.add_shape(1, Inches(M), Inches(1.42 if kicker else 1.06),
                            Inches(W - 2 * M), Emu(14000))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY
    ln.line.fill.background(); ln.shadow.inherit = False
    return s


def footer(s, txt="Source: analyst model.  Strictly Confidential — internal "
                  "circulation only. Not investment advice."):
    tb = s.shapes.add_textbox(Inches(M), Inches(H - 0.46), Inches(W - 2 * M),
                              Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = txt
    r.font.size = Pt(9); r.font.color.rgb = LGREY; r.font.italic = True
    r.font.name = "Calibri"


def bullets(s, items, left=M, top=1.75, width=None, size=15, gap=9):
    width = width or (W - 2 * M)
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                              Inches(H - top - 0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            lead, rest = it
            r = p.add_run(); r.text = "▪  " + lead
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = NAVY
            r.font.name = "Calibri"
            r = p.add_run(); r.text = rest
            r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "▪  " + it
            r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = "Calibri"
    return tb


def tbl(s, headers, rows, left=M, top=1.8, width=None, height=None,
        size=12, hdr_size=12, bold_rows=(), col_w=None):
    width = width or (W - 2 * M)
    nr, nc = len(rows) + 1, len(headers)
    height = height or min(0.36 * nr, H - top - 0.75)
    shp = s.shapes.add_table(nr, nc, Inches(left), Inches(top),
                             Inches(width), Inches(height))
    t = shp.table
    t.first_row = True
    if col_w:
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            t.columns[i].width = Emu(int(Inches(width) * cw / tot))
    for i, htxt in enumerate(headers):
        c = t.cell(0, i)
        c.text = ""
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = c.margin_bottom = Pt(1)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(hdr_size); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = "Calibri"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = t.cell(ri + 1, ci)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = BAND if (ri % 2 == 1) else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Pt(1)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = "Calibri"
            r.font.bold = ri in bold_rows
    return t


def pic(s, png, left, top, width):
    s.shapes.add_picture(png, Inches(left), Inches(top), width=Inches(width))


def statbox(s, left, top, w, h, label, value, sub=None, colour=NAVY):
    box = s.shapes.add_shape(5, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = BAND
    box.line.fill.background(); box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.size = Pt(11); r.font.color.rgb = GREY; r.font.name = "Calibri"
    p = tf.add_paragraph()
    r = p.add_run(); r.text = value
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = colour
    r.font.name = "Calibri"
    if sub:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = sub
        r.font.size = Pt(10); r.font.color.rgb = GREY; r.font.name = "Calibri"


# ------------------------------------------------------------------ 1 title
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(3.1))
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
bg.line.fill.background(); bg.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(M), Inches(0.85), Inches(W - 2 * M), Inches(1.9))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run(); r.text = "JSW ENERGY LIMITED"
r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE
r.font.name = "Calibri"
p = tf.add_paragraph()
r = p.add_run(); r.text = "A thermal balance sheet funding a renewable company"
r.font.size = Pt(22); r.font.italic = True
r.font.color.rgb = RGBColor(0xAF, 0xC8, 0xEA); r.font.name = "Calibri"
p = tf.add_paragraph()
r = p.add_run(); r.text = "Institutional research note  |  July 2026  |  NSE: JSWENERGY"
r.font.size = Pt(14); r.font.color.rgb = RGBColor(0xAF, 0xC8, 0xEA)
r.font.name = "Calibri"
for i, (lab, val, sub, col) in enumerate([
        ("Rating", "HOLD", "derived from the scenario table", NAVY),
        ("CMP (17 Jul 2026)", "Rs.561", "Mkt cap Rs.98,584cr", INK),
        ("Base target price", "Rs.625", "11.9x FY28E EV/EBITDA", NAVY),
        ("Upside", "+11.4%", "bear Rs.286 / bull Rs.839", ORANGE)]):
    statbox(s, M + i * 3.10, 3.55, 2.90, 1.42, lab, val, sub, col)
tb = s.shapes.add_textbox(Inches(M), Inches(5.25), Inches(W - 2 * M), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("The build is real: 13,454MW at FY26A to 25,905MW by FY30E, EBITDA "
          "Rs.10,064cr to Rs.23,536cr. The equity is not. Net debt of Rs.75,038cr "
          "at FY27E-end is 4.7x FY28E EBITDA, so the equity is 29% of enterprise "
          "value and every 1.0x of EV/EBITDA is worth Rs.87 per share — 15.5% of "
          "the price. Reversion to the 16-year mean multiple of 9.0x is Rs.373.")
r.font.size = Pt(14); r.font.color.rgb = INK; r.font.name = "Calibri"
footer(s, "Estimates: attached analyst model, used as given. Primary sources: FY26 "
          "annual report, Q4FY26 and Q1FY27 calls. Market data 17 July 2026. Strictly Confidential — internal circulation only. Not investment advice.")

# ------------------------------------------------------------------ 2 overview
s = slide("Company overview", "India's fifth-largest private generator, mid-transformation")
bullets(s, [
    ("What it is.  ", "5,658MW thermal, 1,631MW hydro and 6,165MW renewables — "
     "13.45GW at FY26A and 14.6GW after Q1FY27, 61% renewable. A further 14GW is "
     "under construction, all tied under long-term PPAs; locked-in capacity is "
     "32.1GW against a 30GW-by-2030 target."),
    ("How it grew.  ", "Almost entirely by acquisition: Mytrah (1,753MW, FY23), "
     "Ind-Barath (700MW, FY24–25), KSK Mahanadi (1,800MW via NCLT, FY25), "
     "O2 Power (4.7GW, Rs.12,468cr EV, FY26). No proprietary technology."),
    ("FY26A was the step-change.  ", "Revenue +61% to Rs.18,901cr, EBITDA +93% to "
     "Rs.10,064cr — almost entirely first-time consolidation of KSK and O2. "
     "Reported PAT rose only 15%, and only because of a Rs.777cr net tax credit."),
    ("What it cost.  ", "Net debt rose Rs.25,155cr in FY26A alone, to Rs.70,081cr. "
     "Leverage 7.0x on the model's measure, 5.2x on the company's (ex-CWIP debt); "
     "interest cover 1.70x."),
    ("The funding is done.  ", "Rs.10,150cr raised across FY26 and Q1FY27 — "
     "Rs.3,000cr promoter preferential, Rs.3,150cr from the JSW Steel stake, "
     "Rs.4,000cr QIP. Leverage improved to 4.95x, inside the 5.0–5.5x guardrail. "
     "Cash Rs.12,881cr."),
    ("Where it trades.  ", "11.2x FY28E EV/EBITDA, 3.2x FY26A book. Its own 12-month "
     "forward multiple has averaged 9.0x since 2010 and sits at 13.9x today."),
], size=15.5, gap=13)
footer(s)

# ------------------------------------------------------------------ 3 structure
s = slide("Corporate structure", "The listed entity produced 18% of FY26A segment EBITDA")
pic(s, "ex7_structure.png", 0.5, 1.62, 8.5)
bullets(s, [
    "The listed parent owns three old thermal plants (2,078MW).",
    "The entire renewable build sits two levels down, inside JSW Neo Energy.",
    "KSK Mahanadi — the largest single EBITDA contributor in FY26A at Rs.3,343cr — "
    "is 74% owned, and the call option on the balance 26% has been served.",
    "Capital sinks: 'Others' lost Rs.642cr of EBITDA in FY26A. The 1GW solar module "
    "plant is on hold; SECI BESS-1 (500MW) is in dispute with no impairment taken.",
    "Entity-level PAT is not disclosed by the company and is not in the model.",
], left=9.25, top=1.75, width=3.5, size=11.5, gap=10)
footer(s, "Source: FY26 Integrated Annual Report (AOC-1), Q4FY26 and Q1FY27 calls, analyst model.")

# ------------------------------------------------------------------ 4 cycles
s = slide("Six business cycles", "Every previous growth cycle ended on a lower return on capital")
tbl(s, ["Cycle", "Revenue (Rs.cr)", "EBITDA margin", "Net debt/EBITDA", "ROCE",
        "What happened"],
    [
        ["1. FY05–08  Captive merchant", "491 → 1,574", "44% → 74%", "2.9x → 1.7x",
         "7% → 29%", "Vijayanagar only; peak merchant tariffs"],
        ["2. FY09–13  First build", "1,591 → 8,934", "18% → 31%", "20.0x → 3.6x",
         "2% → 9%", "Ratnagiri + Barmer; debt to Rs.10,377cr"],
        ["3. FY14–19  Stagnation", "8,705 → 9,138", "37% → 31%", "2.9x → 3.6x",
         "8% → 7%", "No new capacity; Vijayanagar PLF 95% → 50%"],
        ["4. FY20–22  Deleveraging", "8,273 → 8,167", "36% → 44%", "3.3x → 2.2x",
         "9% → 7%", "Cleanest balance sheet in a decade"],
        ["5. FY23–26A  Re-levering", "10,332 → 18,901", "32% → 53%", "6.3x → 7.0x",
         "4% → 6%", "Four acquisitions; net debt up 3.4x"],
        ["6. FY27E–30E  Renewables", "26,535 → 39,506", "49% → 60%", "6.1x → 4.9x",
         "6% → 9%", "12,451MW built; first cycle ROCE rises"],
    ], size=11.5, hdr_size=11.5, col_w=[2.5, 1.7, 1.4, 1.5, 1.1, 3.4], top=1.85)
bullets(s, [
    "ROCE was 29% in FY08, 9% in FY13, 7% in FY22 and 6.4% in FY26A. The model has "
    "it reaching 8.8% by FY30E — still below the 10% WACC in the model's own DCF.",
], top=5.35, size=14)
footer(s)

# ------------------------------------------------------------------ 5 segments
s = slide("Segments", "Renewables go from 40% of EBITDA to 71% by FY30E")
pic(s, "ex2_segment_mix.png", 0.5, 1.7, 7.4)
tbl(s, ["Revenue (Rs.cr)", "FY25A", "FY26A", "FY27E", "FY28E"],
    [[k] + [f"{v:,.0f}" for v in vals[2:]] for k, vals in d.sr.items()
     if k in ("Total Thermal", "Hydro", "Renewables", "Consolidated revenue")],
    left=8.2, top=1.85, width=4.55, size=11.5, hdr_size=11.5, bold_rows=(3,),
    col_w=[2.0, 1.0, 1.0, 1.0, 1.0])
bullets(s, [
    "Thermal is fully built — growth is realisation, not volume.",
    "Hydro converted 97% of revenue to EBITDA in FY26A. No new capacity planned.",
    "Renewables converted 85% vs 49% for thermal — the mix is the margin story.",
    "The model does not carry segment revenue beyond FY28E.",
], left=8.2, top=4.35, width=4.55, size=12, gap=8)
footer(s)

# ------------------------------------------------------------------ 6 model & moat
s = slide("Business model and where the moat sits", "Contracts and cost of capital — not technology")
bullets(s, [
    ("Contractual.  ", "25–40 year PPAs with SECI, NTPC, state discoms; a regulated "
     "15% RoE at Barmer; 14% regulated RoE at Jaigad Transco."),
    ("Site and connectivity.  ", "Land at scale and grid evacuation are the binding "
     "constraint on Indian renewable build — years to assemble."),
    ("Group offtake.  ", "~1,730MW sold to JSW Steel, Cement and Paints. A sticky "
     "customer base and a related-party concentration at the same time."),
    ("Cost of capital.  ", "In a reverse auction everyone buys the same panels. The "
     "cheapest funder wins. AA rating helps; 2.3x net debt/equity does not."),
    ("What there is not.  ", "No technology, no IP, no pricing power. This is an "
     "assets-and-contracts business."),
    ("Capital intensity is the defining feature.  ",
     "Rs.88,885cr gross block on Rs.18,901cr revenue — a 0.25x asset turn, falling "
     "to 0.24x by FY30E. Rs.70,428cr of cumulative FY27E–FY30E operating cash flow "
     "does not cover Rs.89,565cr of capex, let alone Rs.33,098cr of interest."),
], size=15, gap=13)
footer(s)

# ------------------------------------------------------------------ 7 capacity
s = slide("Capacity and capex", "12,451MW to build in four years")
pic(s, "ex3_capacity.png", 0.5, 1.7, 7.6)
tbl(s, ["", "FY26A", "FY27E", "FY28E", "FY29E", "FY30E"],
    [
        ["Capacity (MW)", "13,454", "16,609", "19,996", "22,796", "25,905"],
        ["Added in year (MW)", "2,579", "3,155", "3,387", "2,800", "3,109"],
        ["Storage (MW)", "0", "0", "875", "875", "2,375"],
        ["Capex (Rs.cr)", "10,112", "21,229", "21,864", "21,656", "24,816"],
        ["Capex per MW (Rs.cr)", "3.9", "6.7", "6.5", "7.7", "8.0"],
    ], left=8.35, top=1.85, width=4.4, size=11, hdr_size=11,
    col_w=[1.7, 0.85, 0.85, 0.85, 0.85, 0.85])
bullets(s, [
    "Kutehr hydro: original CoD Sep-2024, commissioned FY26 — 18 months late.",
    "Hetero wind 125MW: CoD 31-Mar-2023 in the tracker, still under construction.",
    "SECI BESS-1 500MW: LoA Jul-2022, in dispute, no impairment taken.",
    "Salboni 1,600MW thermal is under construction but is NOT in the model.",
], left=8.35, top=4.25, width=4.4, size=11.5, gap=8)
footer(s)

# ------------------------------------------------------------------ 8 industry
s = slide("Industry and growth drivers", "Fast supply, slow demand — the tension in the story")
bullets(s, [
    ("Where India is.  ", "533GW installed at Mar-2026; 283.5GW non-fossil (53%), "
     "including 274.7GW renewables. JSW Energy is ~2.7% of national capacity."),
    ("Visible now — renewable share gain.  ", "Solar generation +24% and wind +11% "
     "year-on-year in Q4FY26, against total generation growth of just 3%."),
    ("Visible now — the policy mandate.  ", "500GW of non-fossil by 2030 requires "
     "~54GW of additions a year. The SECI/NTPC/state tender pipeline is the mechanism."),
    ("Prospective — the storage premium.  ", "FDRE, hybrid and RTC tenders price "
     "above plain solar. JSW moves from 451MW of hybrid to 4,798MW plus 3,000MW "
     "BESS and 3,300MW of pumped storage."),
    ("The tension.  ", "Q4FY26 generation growth of 3% was the slowest in six years "
     "while renewable capacity kept compounding. Flat demand plus fast supply is the "
     "classic setup for curtailment and merchant tariff compression."),
], size=15, gap=13)
footer(s, "Source: CEA / MNRE / PIB releases and press reports, July 2026; analyst model.")

# ------------------------------------------------------------------ 9 competition
s = slide("Competitive position", "Cheap capital is the moat — and it is getting more expensive")
tbl(s, ["Peer", "EV/EBITDA", "P/E", "Comment"],
    [
        ["JSW Energy", "11.2x FY28E", "44.0x FY26A",
         "13.5GW; 2.3x net debt/equity; ROCE 6.4%"],
        ["NTPC", "~11.0x", "~12.8x", "~80GW+, sovereign, regulated RoE — the floor"],
        ["Tata Power", "~13.8x", "~33.4x", "Closest structural comparable"],
        ["Adani Power", "~24.4x", "n/a", "Largest listed private thermal by mkt cap"],
        ["Adani Green", ">17x", ">50x", "Pure-play RE — the implicit SOTP anchor"],
    ], top=1.9, size=13, hdr_size=13, bold_rows=(0,),
    col_w=[1.6, 1.3, 1.1, 4.0], height=2.3)
bullets(s, [
    "JSW's ROCE averaged 5.4% over FY23A–FY26A. NTPC on a regulated base and Tata "
    "Power on an integrated model both run materially higher.",
    "A 60% EBITDA margin on a 0.24x asset turn is how a business looks highly "
    "profitable on the income statement and mediocre on capital.",
    "Peer multiples are trailing secondary-source figures, not reconciled to audited "
    "accounts. Like-for-like peer growth, margin and ROCE were not obtainable — a gap.",
], top=4.5, size=13.5, gap=10)
footer(s)

# ------------------------------------------------------------------ 10 thesis
s = slide("Investment thesis", "The market is right about the EBITDA and wrong about who owns it")
bullets(s, [
    ("Variant perception.  ", "The market treats JSW Energy as a renewable "
     "compounder deserving a growth multiple. The EBITDA does more than double. But "
     "92% of that growth is debt-funded renewable capacity, and the equity is 29% of "
     "enterprise value — so the price is a geared bet on the exit multiple."),
    ("1. The build is real and checkable.  ", "13,454MW to 25,905MW; EBITDA "
     "Rs.10,064cr to Rs.23,536cr, a 24% CAGR. 2,579MW came in FY26A, 1,081MW since "
     "April 2026 against a 3,155MW FY27E requirement."),
    ("2. The margin expansion is mechanical.  ", "Renewables convert 85% of revenue "
     "to EBITDA vs 49% for thermal. Consolidated margin 53.2% → 59.6% on mix alone. "
     "Fuel falls from 29.5% of revenue to 25.5%."),
    ("3. But the equity is a thin, geared claim.  ", "Net debt Rs.75,038cr at "
     "FY27E-end = 4.7x FY28E EBITDA. Each 1.0x of multiple = Rs.87/share = 15.5% of "
     "the price. 9.0x gives Rs.373; 11.9x gives Rs.625."),
    ("4. And returns still do not clear the cost of capital.  ",
     "ROCE 6.4% FY26A → 8.8% FY30E, against the 10% WACC in the model's own DCF."),
], size=14.5, gap=12)
footer(s)

# ------------------------------------------------------------------ 11 financials
s = slide("Financial performance", "FY26A doubled EBITDA; PAT grew 15% and only on a tax credit")
tbl(s, ["Rs. cr", "FY26A", "FY27E", "FY28E", "FY29E", "FY30E"],
    [
        ["Revenue"] + [f"{v:,.0f}" for v in d.revenue[4:]],
        ["  growth, yoy"] + [f"{(d.revenue[i]/d.revenue[i-1]-1)*100:,.0f}%"
                             for i in range(4, 9)],
        ["EBITDA"] + [f"{v:,.0f}" for v in d.ebitda[4:]],
        ["  EBITDA margin"] + [f"{d.ebitda[i]/d.revenue[i]*100:.1f}%"
                               for i in range(4, 9)],
        ["Interest expense"] + [f"{v:,.0f}" for v in d.interest[4:]],
        ["Effective tax rate"] + [f"{v:.1f}%" for v in d.taxrate[4:]],
        ["PAT after minority"] + [f"{v:,.0f}" for v in d.patrep[4:]],
        ["EPS (Rs.)"] + [f"{v:.2f}" for v in d.eps[4:]],
        ["ROE"] + [f"{v:.1f}%" for v in d.roe[4:]],
        ["ROCE"] + [f"{v:.1f}%" for v in d.roce[4:]],
    ], top=1.85, size=12.5, hdr_size=12.5, bold_rows=(0, 2, 6),
    col_w=[2.2, 1.0, 1.0, 1.0, 1.0, 1.0], width=7.4, height=3.9)
pic(s, "ex1_revenue_margin.png", 8.15, 1.75, 4.6)
bullets(s, [
    "FY26A's tax charge was NEGATIVE Rs.777cr on positive PBT — deferred tax "
    "assets recognised at Utkal and KSK once PPAs were signed. Explained and "
    "legitimate, but not repeatable: management guides 23–24% from here. "
    "Normalised, FY26A PAT is ~Rs.969cr, not Rs.2,239cr — +15% becomes about –50%.",
], top=5.95, size=13.5, width=7.4)
footer(s)

# ------------------------------------------------------------------ 12 cost & returns
s = slide("Cost structure and returns", "Margin is mix; cash conversion is not the problem — capex is")
tbl(s, ["% of revenue", "FY26A", "FY27E", "FY28E", "FY29E", "FY30E"],
    [
        ["Fuel cost"] + [f"{d.fuel[i]/d.revenue[i]*100:.1f}%" for i in range(4, 9)],
        ["Employee cost"] + [f"{d.staff[i]/d.revenue[i]*100:.1f}%"
                             for i in range(4, 9)],
        ["Other expenses"] + [f"{d.otherexp[i]/d.revenue[i]*100:.1f}%"
                              for i in range(4, 9)],
        ["EBITDA margin"] + [f"{d.ebitda[i]/d.revenue[i]*100:.1f}%"
                             for i in range(4, 9)],
    ], top=1.85, size=11.5, hdr_size=11.5, bold_rows=(3,),
    col_w=[1.9, 0.86, 0.86, 0.86, 0.86, 0.86], width=6.55, height=1.7)
tbl(s, ["Rs. cr", "FY26A", "FY27E", "FY28E", "FY29E", "FY30E"],
    [
        ["Operating cash flow"] + [f"{v:,.0f}" for v in d.ocf[4:]],
        ["Capex"] + [f"({v:,.0f})" for v in d.capex[4:]],
        ["Free cash flow"] + [f"({abs(v):,.0f})" for v in d.fcf[4:]],
        ["FCF post interest"] + [f"({abs(v):,.0f})" for v in d.fcf_int[4:]],
        ["Pre-tax OCF/EBITDA"] + [f"{v:.0f}%" for v in d.ocf_ebitda[4:]],
    ], top=4.1, size=11.5, hdr_size=11.5, bold_rows=(3,),
    col_w=[1.62, 0.99, 0.99, 0.99, 0.99, 0.99], width=6.55, height=2.05)
bullets(s, [
    ("Fuel is the only variable cost.  ", "Renewables and hydro have none, so the "
     "cost base goes from ~63% fixed at the FY26A mix to ~75% fixed by FY30E."),
    ("Operating leverage is asymmetric.  ", "A 10% renewable generation shortfall "
     "hits EBITDA almost one-for-one; a 10% thermal shortfall costs about half that."),
    ("Cash converts at the EBITDA line.  ", "Pre-tax OCF/EBITDA of 101–110% through "
     "the forecast; debtor days improved from 81 to 57 in FY26A."),
    ("It does not convert at the equity line.  ", "Cumulative FCF after interest of "
     "NEGATIVE Rs.53,993cr over FY27E–FY30E. There is no year in the model in which "
     "the company self-funds."),
], left=7.35, top=1.85, width=5.4, size=12.5, gap=11)
footer(s)

# ------------------------------------------------------------------ 13 balance sheet
s = slide("Balance sheet and capital allocation", "Rs.73,000cr deployed in five years; ROCE went from 7.1% to 6.4%")
pic(s, "ex5_leverage.png", 0.5, 1.75, 6.6)
bullets(s, [
    ("The record, FY22A–FY26A.  ", "Rs.25,007cr of operating cash flow. Spent "
     "Rs.31,382cr on capex and Rs.26,537cr acquiring subsidiaries, paid Rs.11,760cr "
     "of interest and ~Rs.1,700cr of dividends. Funded with Rs.66,953cr of new debt "
     "and Rs.6,069cr of new equity."),
    ("The path from here.  ", "Net debt rises every single year, Rs.70,081cr to "
     "Rs.115,314cr. Leverage falls from 7.0x to 4.9x only because EBITDA rises "
     "faster. Management's sub-5.0x target is reached in FY29E at the earliest."),
    ("The plan is funded.  ", "Rs.10,150cr raised; Rs.24,184cr of capex is now "
     "contractually committed in the audited accounts; reported cost of debt fell "
     "67bp to 8.36%. The CFO says FY27's Rs.20,000cr needs no further equity."),
    ("Two flags remain.  ", "A dividend maintained through four years of large "
     "negative free cash flow; and Rs.7,862cr of JSW Steel shares still held while "
     "borrowing at 8.36% — a cross-holding that runs both ways, since JSW Steel "
     "owns 4.86% of JSW Energy."),
], left=7.35, top=1.85, width=5.4, size=12.5, gap=11)
footer(s)

# ------------------------------------------------------------------ 14 price story
s = slide("The multiple, 2010–2026", "EBITDA doubled in FY26A and the multiple compressed")
pic(s, "ex6_evband.png", 0.5, 1.7, 7.6)
bullets(s, [
    ("FY11–20.  ", "Derated 9–11x to 5–6x, trough 3.9x. No new capacity; "
     "Vijayanagar PLF 95% → 50%."),
    ("FY21–22.  ", "5.8x to 22.4x. JSW Neo formed, 20GW target announced, leverage "
     "at 2.2x — the cleanest balance sheet in a decade."),
    ("FY23.  ", "Back to 10.3x. EBITDA fell 8% on imported coal; leverage jumped to "
     "6.3x on Mytrah."),
    ("FY24–25.  ", "Back to a 20.1x peak on Ind-Barath, KSK and a Rs.4,944cr QIP."),
    ("FY26A to date.  ", "12.0x–14.9x, now 13.9x. EBITDA nearly doubled while the "
     "average multiple fell from 16.9x (2024) to 13.2x (2026). The market took the "
     "EBITDA and refused to keep paying for it."),
    ("Gap.  ", "A monthly closing-price series is not in the model — the narrative "
     "is told through the multiple plus the 52-week range of Rs.428–617."),
], left=8.25, top=1.8, width=4.5, size=11.5, gap=9)
footer(s)

# ------------------------------------------------------------------ 15 valuation
s = slide("Valuation", "Base Rs.625 = 11.9x FY28E EV/EBITDA less FY27E-end net debt")
tbl(s, ["Scenario", "Mult.", "Applied to", "EBITDA", "Equity value", "Target",
        "vs CMP"],
    [
        ["Bear", "8.0x", "FY28E", "15,942", "52,497", "Rs.286", "(49%)"],
        ["Downside case", "10.2x", "FY28E", "15,942", "87,569", "Rs.478", "(15%)"],
        ["Base", "11.9x", "FY28E", "15,942", "114,671", "Rs.625", "+11%"],
        ["Bull", "11.9x", "FY29E", "20,212", "153,849", "Rs.839", "+50%"],
        ["16-yr mean multiple", "9.0x", "FY28E", "15,942", "68,439", "Rs.373", "(33%)"],
    ], top=1.85, size=12, hdr_size=12, bold_rows=(2,),
    col_w=[2.3, 0.85, 1.1, 1.05, 1.3, 1.0, 0.95], width=8.6, height=2.2)
bullets(s, [
    ("Sensitivity.  ", "FY28E EBITDA 10% below forecast → Rs.522, a 7% DOWNSIDE. "
     "Margin at 50.0% instead of 52.5% → Rs.576. Mean-reversion to 9.0x → Rs.373."),
    ("Minority adjustment.  ", "The call option on KSK's 26% has been served, so "
     "full consolidation is right — but the ~Rs.1,170cr price is not in FY27E net "
     "debt. Adjusting takes the base to ~Rs.619, not the ~Rs.582 it would have been "
     "without the buyout."),
    ("Not credited in the base.  ", "Residual JSW Steel stake Rs.34/share; Salboni 3,200MW; "
     "KSK units 3–6."),
    ("The cross-check.  ", "The model's own asset-based SOTP is Rs.510 — BELOW the "
     "market price. The two methods bracket the CMP. That is what fairly valued "
     "looks like."),
    ("Why no DCF.  ", "FCFE is negative in four of seven forecast years and the "
     "terminal value is negative; the model's DCF returns minus Rs.376 per share. "
     "An artefact, not a signal. No weight placed on it."),
], top=4.35, size=12.5, gap=8)
footer(s)

# ------------------------------------------------------------------ 16 governance
s = slide("Governance and management", "Clean audit, aligned promoter, one open item")
tbl(s, ["Parameter", "Status", "Comment"],
    [
        ["Promoter holding", "66.5% (Jun-26) from 69.3% (Sep-24)",
         "Dilution from QIPs, not selling. Rs.625cr of warrants already subscribed "
         "at Rs.525; Rs.1,875cr still payable."],
        ["Institutional trend", "FII 14.9% → 11.4%; DII 9.8% → 16.2%",
         "Foreign investors have sold for two years; domestic institutions absorbed "
         "the stock and the QIP."],
        ["Concalls / disclosure", "Quarterly, granular",
         "Segment and commissioning detail above sector norm."],
        ["Related party", "Material, quantum unverified",
         "~1,730MW sold to JSW Steel, Cement, Paints; fuel JV with a state entity."],
        ["Auditor", "Deloitte Haskins & Sells LLP — unmodified",
         "Clean FY26 opinion, no emphasis of matter. One key audit matter: tariff "
         "disputes with customers."],
        ["Contingent liabilities", "~Rs.4,506cr, ~15% of net worth",
         "Claims Rs.2,520cr (down yoy), tariff disputes Rs.340cr, guarantees "
         "Rs.605cr, JV share Rs.1,041cr. Not growing."],
        ["Promoter pledge", "Not disclosed in the AR",
         "Only project-level lender pledges. Still to be checked in the quarterly "
         "shareholding filing — the one open item."],
        ["Guidance vs delivery", "On plan in aggregate, late on projects",
         "Kutehr 18 months late; Hetero wind 3 years late; SECI BESS-1 stalled; "
         "solar module plant on hold. But 2,579MW did commission in FY26A."],
    ], top=1.85, size=12, hdr_size=12, col_w=[2.4, 2.6, 5.1], height=3.6)
bullets(s, [
    "Sajjan Jindal (Chairman & MD) — the key-man and capital-allocation authority. "
    "Sharad Mahendra (Joint MD & CEO since Feb-2024). Chandrasekaran Prabhakaran "
    "(CFO since 1 Jan 2026, replacing Pritesh Vinay) — a CFO change six months "
    "before the sector's largest capital raise, which was then delivered.",
    "The model's own target of 28.3GW by FY30E is ~6% below the company's stated "
    "30GW — the estimates used here are slightly more conservative than guidance.",
], top=5.7, size=13, gap=8)
footer(s)

# ------------------------------------------------------------------ 17 triggers & risks
s = slide("Triggers and risks", "What to watch, and what would break the thesis")
tb = s.shapes.add_textbox(Inches(M), Inches(1.75), Inches(6.1), Inches(4.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "TRIGGERS — dated and checkable"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
for t in ["Q2FY27 (Oct-2026): EBITDA must run Rs.3,360cr/quarter for 9 months vs "
          "Rs.2,873cr in 1QFY27.",
          "Standalone segment revenue: FY27E assumes Rs.6,612cr vs Rs.3,030cr in "
          "FY26A and Rs.1,101cr delivered in 1QFY27.",
          "KSK minority buyout price — not yet crystallised, expected end-Q2FY27; "
          "~Rs.1,170cr is absent from the model's net debt.",
          "First 600MW of KSK units 3-6, mid-FY27 — pure upside, ~Rs.35/share.",
          "Rajasthan evacuation line, July 2026 — ends the curtailment that cost "
          "Rs.50cr in FY26A.",
          "Quarterly MW commissioned — 3,155MW needed in FY27E, 873MW came in Q1.",
          "Final DSM regulations — 1.5-2% of RE revenue on management's worst case; "
          "substation-level grouping would reduce it."]:
    p = tf.add_paragraph(); p.space_after = Pt(7)
    r = p.add_run(); r.text = "▪  " + t
    r.font.size = Pt(12.5); r.font.color.rgb = INK; r.font.name = "Calibri"

tb = s.shapes.add_textbox(Inches(6.95), Inches(1.75), Inches(5.8), Inches(4.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "RISKS — most damaging first"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ORANGE
r.font.name = "Calibri"
for t in ["Multiple compression. Each 1.0x = Rs.87/share. Mean reversion to 9.0x is "
          "Rs.373, a 33% loss with no operational change.",
          "Renewable execution slippage. A one-year slip removes Rs.3,000–4,000cr of "
          "FY30E EBITDA — over Rs.200/share of target.",
          "Funding. De-risked for FY27 (raise complete, cash Rs.12,881cr); "
          "unresolved for FY29-30, ~Rs.46,000cr on 2.1x interest cover.",
          "Evacuation, not merchant price. Only 9,500 of a planned 15,000 ckm built; "
          "eases ~2029. Open capacity is just ~5% of the base.",
          "Three items the model omits: the non-repeatable FY26A tax credit, the "
          "cost of the KSK minority buyout, and the Supreme Court's 18% Himachal "
          "free-power ruling (~287 MU/yr, permanent).",
          "Per-MW economics. The model implies Rs.1.02cr/MW of incremental RE "
          "EBITDA against management's own Rs.0.75cr/MW steady state — a "
          "Rs.3,325cr gap to FY30E."]:
    p = tf.add_paragraph(); p.space_after = Pt(7)
    r = p.add_run(); r.text = "▪  " + t
    r.font.size = Pt(12.5); r.font.color.rgb = INK; r.font.name = "Calibri"

box = s.shapes.add_shape(5, Inches(M), Inches(6.35), Inches(W - 2 * M), Inches(0.72))
box.fill.solid(); box.fill.fore_color.rgb = BAND
box.line.fill.background(); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = Pt(10)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Falsification test:  "
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = ORANGE
r.font.name = "Calibri"
r = p.add_run()
r.text = ("sell if FY27 EBITDA lands below Rs.11,500cr AND net debt/EBITDA is above "
          "6.5x at FY27-end — the build has neither delivered nor been funded. "
          "Upgrade to buy if FY27 EBITDA exceeds Rs.13,000cr and leverage closes "
          "below 5.5x on the model's all-in measure.")
r.font.size = Pt(12.5); r.font.color.rgb = INK; r.font.name = "Calibri"
footer(s)

# ------------------------------------------------------------------ 18 outlook
s = slide("Final outlook", "HOLD  |  Base TP Rs.625  |  +11.4%  |  CMP Rs.561")
for i, (lab, val, sub, col) in enumerate([
        ("Bear  8.0x FY28E", "Rs.286", "(49%)", ORANGE),
        ("Base  11.9x FY28E", "Rs.625", "+11%", NAVY),
        ("Bull  11.9x FY29E", "Rs.839", "+50%", BLUE),
        ("Model SOTP", "Rs.510", "(9%) — below CMP", GREY)]):
    statbox(s, M + i * 3.10, 1.85, 2.90, 1.4, lab, val, sub, col)
bullets(s, [
    ("The business is genuinely better than it was.  ",
     "From 6,605MW to 13,454MW in four years, from a stagnant thermal generator to "
     "a company where renewables supply 71% of EBITDA by FY30E, on 25–40 year "
     "fixed-tariff PPAs, with an AA rating and a promoter subscribing equity at Rs.525."),
    ("The base target requires nothing heroic.  ",
     "Only that FY28E EBITDA of Rs.15,942cr is delivered and the market keeps paying "
     "roughly what it pays today. No re-rating, no acquisition, and no credit for "
     "Salboni, KSK units 3–6 or the Rs.34/share residual JSW Steel stake."),
    ("But the risk-reward is not asymmetric.  ",
     "11% base upside against 33% downside on a mean-reverting multiple; a 10% "
     "EBITDA miss alone flips +11% to –7%; and terminal ROCE of 8.8% is still below "
     "the model's own 10% WACC. An asset-based SOTP of Rs.510 and a forward-multiple "
     "target of Rs.625 bracket the current price."),
    ("What caps position sizing.  ",
     "Less than we thought before reading the primary sources — the audit is clean "
     "and contingent exposure is ~15% of net worth and not growing. What remains: "
     "~Rs.46,000cr of FY29-30 capex on 2.1x interest cover, the Himachal free-power "
     "ruling absent from the forecast, per-MW economics 36% above management's own "
     "benchmark, and an unseen promoter pledge position. Buy on evidence — FY27 "
     "EBITDA above Rs.13,000cr with leverage under 5.5x — or nearer Rs.480."),
], top=3.45, size=12.6, gap=8)
footer(s)

prs.save("JSW_Energy_Pitch_Deck_July2026.pptx")
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
